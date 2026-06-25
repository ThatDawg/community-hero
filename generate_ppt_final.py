from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

template_path = r'C:\Users\Lenovo uSER\Downloads\SOCF PPT.pptx'
prs = Presentation(template_path)

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
ACCENT2 = RGBColor(0x00, 0xD2, 0xFF)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
RED = RGBColor(0xFF, 0x55, 0x55)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

def get_ph(slide, idx):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None

def set_line_spacing_1_5(shape):
    """Set 1.5 line spacing on all paragraphs in a shape"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        para.line_spacing = 1.5

def set_title(slide, text):
    title = get_ph(slide, 0)
    if title:
        title.text = ""
        p = title.text_frame.paragraphs[0]
        p.text = text
        p.line_spacing = 1.5
        for run in p.runs:
            run.font.color.rgb = WHITE
            run.font.bold = True
            run.font.size = Pt(28)

def set_content(slide, items):
    content = get_ph(slide, 1)
    if not content:
        return
    tf = content.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.line_spacing = 1.5
        p.space_after = Pt(4)

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def tb_set(tf, text, size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.line_spacing = 1.5
    p.alignment = alignment

def tb_add_para(tf, text, size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.line_spacing = 1.5
    p.alignment = alignment
    p.space_before = Pt(2)
    return p

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

# ======================================================================
# SLIDE 1 — TITLE
# ======================================================================
slide = prs.slides[0]
set_title(slide, "Personal Study Planner Agent")
subtitle = get_ph(slide, 1)
if subtitle:
    p = subtitle.text_frame.paragraphs[0]
    p.text = "Proposed by : [Your Team Name]"
    p.font.color.rgb = WHITE
    p.font.size = Pt(18)
    p.font.bold = False
    p.line_spacing = 1.5

# ======================================================================
# SLIDE 2 — INSTRUCTIONS
# ======================================================================
slide = prs.slides[1]
set_title(slide, "Instructions")
set_content(slide, [
    "Keep the presentation within the prescribed slide limit.",
    "Allowed Font: Century Gothic, Times New Roman, Arial",
    "Use concise text, visuals, and diagrams wherever appropriate.",
    "Clearly explain the problem, proposed solution, feasibility, and expected impact.",
    "Be realistic about challenges, risks, and implementation constraints.",
    "",
    "Problem Statement: Personal Study Planner Agent (Track 2 \u2014 AI Agent)",
    "Team: [Your Team Name]  |  Round 1 \u2014 Ideathon Submission",
])

# ======================================================================
# SLIDE 3 — PROPOSED SOLUTION
# ======================================================================
slide = prs.slides[2]
set_title(slide, "Proposed Solution")
set_content(slide, [
    "An AI Agent that transforms scattered academic inputs (syllabi, timetable, assignments) into personalized weekly study plans",
    "Students upload PDFs or enter deadlines manually; the agent parses, prioritizes, and generates an adaptive calendar",
    "Key Features: Auto-plan generation, interactive editing (accept/edit/reject), progress tracking, smart reminders, transparent reasoning",
    "Target Users: College students managing 5\u20136 courses, hostel students with fixed study hours, part-time learners",
    "Workflow: Upload Inputs \u2192 Parse via RAG \u2192 LLM Plan Generation \u2192 Review & Edit \u2192 Adaptive Tracking",
])

# ======================================================================
# SLIDE 4 — TECHNICAL STRATEGY (Tech/Tools)
# ======================================================================
slide = prs.slides[3]
set_title(slide, "Technical Strategy")
set_content(slide, [
    "Technologies / Tools:",
    "",
    "1. Languages: Python (primary)",
    "   \u2022 Frontend: Streamlit (MVP) / React (production)",
    "   \u2022 Backend: FastAPI",
    "   \u2022 AI/ML: GPT-4 / Gemini + LangChain + LangGraph",
    "   \u2022 Database: SQLite (dev) / PostgreSQL (production)",
    "   \u2022 Vector DB: ChromaDB (RAG for syllabus retrieval)",
    "   \u2022 Other: PyMuPDF (PDF parser), Docker",
    "",
    "2. Brief Description: End-to-end agent using LLM reasoning with RAG. Modular design allows hot-swapping models. Local-first with cloud option.",
])

# ======================================================================
# SLIDE 5 — TECHNICAL STRATEGY (Methodology)
# ======================================================================
slide = prs.slides[4]
set_title(slide, "Technical Strategy")
set_content(slide, [
    "Methodology",
    "",
    "1. Ingestion: Parse PDF syllabi using PyMuPDF + regex. Store extracted courses, deadlines, exams in ChromaDB vector database for RAG.",
    "",
    "2. Planning: LangChain agent queries LLM with parsed context, deadline proximity, task difficulty, and user preferences to generate ranked weekly plans.",
    "",
    "3. Conflict Detection: Rule-based engine checks for overlapping deadlines, unbalanced workload, and prerequisite ordering before finalization.",
    "",
    "4. Feedback Loop: User accepts, edits, or rejects suggestions. Preferences stored to personalize future planning iterations. Agent explains reasoning on demand.",
])

# ======================================================================
# SLIDE 6 — TECHNICAL STRATEGY (Flowchart)
# ======================================================================
slide = prs.slides[5]
set_title(slide, "Technical Strategy")

# Remove old content placeholder from slide 6
# Content is now in a flowchart below

box_w = Inches(8.5)
box_h = Inches(0.55)
arrow_h = Inches(0.25)
start_y = Inches(2.3)
x_center = Inches(2.4)

flow_items = [
    ("INPUT: Syllabus PDF / Timetable / Assignment List", ACCENT),
    ("PARSER: PyMuPDF + Regex \u2192 Extract courses, exams, deadlines", GREEN),
    ("VECTOR DB: ChromaDB \u2192 Store course context for RAG", ORANGE),
    ("LLM AGENT: LangChain + GPT-4/Gemini \u2192 Generate weekly plan", ACCENT),
    ("VALIDATOR: Conflict Detection \u2192 Overlaps, dependencies, balance", RED),
    ("UI: Streamlit Calendar \u2192 Accept / Edit / Reject plan", ACCENT2),
    ("TRACKER: Progress \u2192 Mark done \u2192 Auto-adaptive replanning", GREEN),
]

for i, (text, color) in enumerate(flow_items):
    y = start_y + i * (box_h + arrow_h + Inches(0.05))
    box = add_rect(slide, x_center, y, box_w, box_h, color)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.line_spacing = 1.5
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if i < len(flow_items) - 1:
        arrow_y = y + box_h
        add_arrow(slide, x_center + box_w // 2 - Inches(0.15), arrow_y, Inches(0.3), arrow_h)

# ======================================================================
# SLIDE 7 — FEASIBILITY ANALYSIS
# ======================================================================
slide = prs.slides[6]
set_title(slide, "Feasibility and Viability")
set_content(slide, [
    "Feasibility Analysis",
    "",
    "Can the solution be realistically built?\nYes. All components (PDF parsing, LLM APIs, Streamlit, ChromaDB, LangChain) are mature, well-documented, and available with free or low-cost tiers.",
    "",
    "What resources are required?\n2\u20133 developers familiar with Python. ~$20 in LLM API credits for dev. Standard laptop. No specialized hardware needed.",
    "",
    "Why is it practical and achievable?\nModular architecture enables parallel development. MVP achievable in 4 weeks. Core logic works with offline fallback models (Llama 3 via Ollama).",
])

# ======================================================================
# SLIDE 8 — STRATEGIES TO OVERCOME CHALLENGES
# ======================================================================
slide = prs.slides[7]
set_title(slide, "Feasibility and Viability")
set_content(slide, [
    "Strategies To Overcome Challenges",
    "",
    "Challenge 1: LLM hallucination in plan generation\n\u2192 RAG grounding with syllabus text in vector DB, strict JSON output formatting, mandatory human review.",
    "",
    "Challenge 2: Inconsistent PDF parsing across formats\n\u2192 Template-based extraction with regex. Manual entry fallback for unsupported formats.",
    "",
    "Challenge 3: User adoption friction\n\u2192 Minimal onboarding (upload & go), intuitive drag-drop calendar, instant value demonstration.",
    "",
    "Backup Plan: Rule-based heuristic planner activates with predefined scheduling rules if LLM API is unavailable.",
])

# ======================================================================
# SLIDE 9 — IMPACT AND BENEFITS
# ======================================================================
slide = prs.slides[8]
set_title(slide, "Impact and Benefits Envisaged")
set_content(slide, [
    "Expected Impact:",
    "  \u2022 Reduces weekly planning time from ~2 hours to under 5 minutes (95% reduction)",
    "  \u2022 Increases deadline adherence by ~40% through automated reminders and conflict alerts",
    "  \u2022 Eliminates scheduling conflicts with real-time dependency and overlap checking",
    "",
    "Key Benefits:",
    "  \u2022 Students: Less stress, better grades through distributed study, improved time management",
    "  \u2022 Institutions: Better academic metrics, reduced dropout risk, workload pattern insights",
    "",
    "Long-term Value:",
    "  \u2022 Scalable to entire university; integrates with LMS (Moodle, Canvas, Blackboard) via API",
    "  \u2022 Anonymized data enables curriculum optimization based on actual student workload patterns",
])

# ======================================================================
# SLIDE 10 — WHY THIS IDEA MAY FAIL
# ======================================================================
slide = prs.slides[9]
set_title(slide, "Why This Idea May Fail ?")
set_content(slide, [
    "1. LLM Hallucination Risk: Unrealistic plans from inaccurate parsing or invented deadlines.",
    "   \u2192 Mitigation: RAG grounding, strict validation rules, mandatory user review before acceptance.",
    "",
    "2. Low Adoption: Students may not trust AI plans or find data entry tedious.",
    "   \u2192 Mitigation: PDF upload with auto-parsing, instant value on first use, gamified progress tracking.",
    "",
    "3. Data Privacy: Concern about uploading materials to external LLM APIs.",
    "   \u2192 Mitigation: Local-first architecture. Option for offline models (Llama 3, Mistral) via Ollama.",
    "",
    "4. Scope Creep: Too many features dilute core planning.",
    "   \u2192 Strict MVP scope: planning, editing, tracking only. Bonus features strictly post-MVP.",
])

# ======================================================================
# SLIDE 11 — REFERENCES
# ======================================================================
slide = prs.slides[10]
set_title(slide, "References")
set_content(slide, [
    "1. Research Papers:",
    '   \u2022 Cred\u00e9, M., & Kuncel, N. R. (2008). "Study Habits, Skills, and Attitudes." Perspectives on Psychological Science.',
    '   \u2022 Pintrich, P. R. (2004). "Conceptual Framework for Motivation and Self-Regulated Learning." Educational Psychology Review.',
    '   \u2022 Zimmerman, B. J. (2002). "Becoming a Self-Regulated Learner." Theory Into Practice.',
    "",
    "2. Tools & Frameworks:",
    "   \u2022 LangChain \u2014 https://www.langchain.com/",
    "   \u2022 ChromaDB \u2014 https://www.trychroma.com/",
    "   \u2022 Streamlit \u2014 https://streamlit.io/",
    "   \u2022 FastAPI \u2014 https://fastapi.tiangolo.com/",
    "   \u2022 PyMuPDF \u2014 https://pymupdf.readthedocs.io/",
])

# ======================================================================
# SLIDE 12 — THANK YOU
# ======================================================================
slide = prs.slides[11]
set_title(slide, "Thank You")
subtitle = get_ph(slide, 1)
if subtitle:
    tf = subtitle.text_frame
    for p in tf.paragraphs:
        p.clear()
    p1 = tf.paragraphs[0]
    p1.text = "Team [Your Team Name]  |  Track 2 \u2014 AI Agent"
    p1.font.color.rgb = WHITE
    p1.font.size = Pt(18)
    p1.line_spacing = 1.5

    p2 = tf.add_paragraph()
    p2.text = "Personal Study Planner Agent  |  GSoC Innovators' Club \u2014 Round 1"
    p2.font.color.rgb = WHITE
    p2.font.size = Pt(16)
    p2.line_spacing = 1.5

    p3 = tf.add_paragraph()
    p3.text = "[Email]  |  [GitHub]  |  [LinkedIn]"
    p3.font.color.rgb = LIGHT_GRAY
    p3.font.size = Pt(14)
    p3.line_spacing = 1.5

# ======================================================================
# APPLY 1.5 LINE SPACING to ALL slides (any remaining shapes)
# ======================================================================
for slide in prs.slides:
    for shape in slide.shapes:
        set_line_spacing_1_5(shape)

# ======================================================================
# SAVE
# ======================================================================
output_path = r'C:\Users\Lenovo uSER\Desktop\Project1\Personal_Study_Planner_Agent_Round1.pptx'
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print("All line spacing set to 1.5. Flowchart added on slide 6.")
